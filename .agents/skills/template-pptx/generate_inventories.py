"""Generate layout_inventory.txt and shape_inventory.txt for a template.

Usage:
    python generate_inventories.py <template.pptx> [output_dir]

If output_dir is omitted, files are written next to the .pptx file.
"""

import sys
from pathlib import Path
from pptx import Presentation


def _shape_line(j: int, shape) -> str:
    """Format a single shape as one compact line."""
    name = shape.name

    # Type label — placeholder type for placeholders, empty for others
    type_label = ""
    if shape.is_placeholder:
        type_label = str(shape.placeholder_format.type)

    # Text preview (pipe-separated paragraphs, truncated)
    text = ""
    if shape.has_text_frame:
        parts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
        text = " | ".join(parts)
        text = text.replace("\n", " ").replace("\r", " ").replace("\x0b", " ")
        if len(text) > 100:
            text = text[:97] + "..."

    pieces = [f"{j:02d}: {name}"]
    if type_label:
        pieces.append(type_label)
    if text:
        pieces.append(text)

    return "  ".join(pieces)


def layout_inventory(prs: Presentation) -> str:
    lines = []
    w = prs.slide_width / 914400
    h = prs.slide_height / 914400
    lines.append(f"Slide size: {w:.2f}\" x {h:.2f}\"")
    lines.append("")

    for i, layout in enumerate(prs.slide_layouts):
        lines.append(f"--- LAYOUT {i}: {layout.name} ---")

        # All shapes on the layout (placeholders + decorative/branding elements)
        all_layout_shapes = list(layout.shapes)

        for j, shape in enumerate(all_layout_shapes, 1):
            lines.append(_shape_line(j, shape))

        if not all_layout_shapes:
            lines.append("  (empty layout)")
        lines.append("")

    return "\n".join(lines)


def shape_inventory(prs: Presentation) -> str:
    lines = []

    for slide_num, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        lines.append(f"--- SLIDE {slide_num} ({layout_name}) ---")

        for j, shape in enumerate(slide.shapes, 1):
            lines.append(_shape_line(j, shape))

        lines.append("")

    return "\n".join(lines)


def main():
    if len(sys.argv) < 2:
        print(__doc__.strip())
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else pptx_path.parent

    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(pptx_path))

    layout_path = output_dir / "layout_inventory.txt"
    layout_path.write_text(layout_inventory(prs), encoding="utf-8")
    print(f"Wrote {layout_path}")

    shape_path = output_dir / "shape_inventory.txt"
    shape_path.write_text(shape_inventory(prs), encoding="utf-8")
    print(f"Wrote {shape_path}")


if __name__ == "__main__":
    main()
